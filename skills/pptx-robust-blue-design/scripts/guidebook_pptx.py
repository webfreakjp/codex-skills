#!/usr/bin/env python3
"""Create or restyle PPTX decks in a robust blue guidebook-like visual style."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as exc:
    raise SystemExit("Missing dependency: install with `python -m pip install --user python-pptx`.") from exc


BLUE = RGBColor(0, 49, 216)
BLUE_2 = RGBColor(38, 74, 244)
PALE_BLUE = RGBColor(232, 241, 254)
TEXT = RGBColor(26, 26, 28)
GRAY = RGBColor(148, 148, 148)
LIGHT_GRAY = RGBColor(230, 230, 230)
MID_GRAY = RGBColor(98, 98, 100)
WHITE = RGBColor(255, 255, 255)
CARD_BORDER = RGBColor(190, 205, 235)
FONT_JP = "Yu Gothic"


def text_of_slide(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                texts.append(text)
    return texts


def clear_slides(prs: Presentation) -> None:
    while len(prs.slides):
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]


def set_run(run, size: float, color: RGBColor, bold: bool = False) -> None:
    run.font.name = FONT_JP
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def textbox(slide, left, top, width, height, text: str, size: float, color: RGBColor,
            bold: bool = False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = anchor
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = str(text or "")
    set_run(run, size, color, bold)
    return box


def rect(slide, left, top, width, height, fill: RGBColor, line: RGBColor | None = None,
         rounded: bool = False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.8)
    return shape


def add_heading(slide, title: str) -> None:
    textbox(slide, Inches(0.55), Inches(0.52), Inches(7.2), Inches(0.5), title, 23, BLUE, True)
    rect(slide, Inches(0.55), Inches(1.17), Inches(4.2), Inches(0.012), LIGHT_GRAY)
    rect(slide, Inches(0.55), Inches(1.17), Inches(1.0), Inches(0.025), BLUE)


def add_page_number(slide, page_no: int, color: RGBColor = GRAY) -> None:
    textbox(slide, Inches(11.35), Inches(6.88), Inches(1.05), Inches(0.12),
            str(page_no), 6, color, align=PP_ALIGN.RIGHT)


def add_footer(slide, page_no: int, steps: list[str] | None = None, active: int | None = None) -> None:
    labels = steps or []
    labels = labels[:5]
    if not labels:
        add_page_number(slide, page_no)
        return
    left = Inches(0.55)
    top = Inches(6.25)
    gap = Inches(0.05)
    width = Inches(11.5 / len(labels)) - gap
    for idx, label in enumerate(labels):
        color = BLUE_2 if active == idx else LIGHT_GRAY
        x = left + idx * (width + gap)
        rect(slide, x, top, width, Inches(0.16), color)
        textbox(slide, x, top - Inches(0.18), width, Inches(0.12), label, 5.5,
                BLUE if active == idx else GRAY, active == idx, PP_ALIGN.CENTER)
    add_page_number(slide, page_no)


def add_optional_progress(slide, page_no: int, data: dict[str, Any]) -> None:
    if data.get("show_progress"):
        add_footer(slide, page_no, data.get("steps"), data.get("active"))
    else:
        add_page_number(slide, page_no)


def bullet_text(value: Any) -> str:
    if isinstance(value, list):
        return "\n".join(f"・{item}" for item in value)
    return str(value or "")


def small_label(slide, left, top, text: str, color: RGBColor = BLUE) -> None:
    textbox(slide, left, top, Inches(2.0), Inches(0.14), text.upper(), 7.5, color, True)


def cover_slide(prs: Presentation, data: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, prs.slide_width, prs.slide_height, BLUE)
    rect(slide, Inches(0.72), Inches(1.36), Inches(0.09), Inches(2.35), WHITE)
    textbox(slide, Inches(0.98), Inches(1.42), Inches(8.65), Inches(1.55),
            data.get("title", "Untitled"), 35, WHITE, True)
    if data.get("subtitle"):
        textbox(slide, Inches(1.00), Inches(3.28), Inches(6.9), Inches(0.38),
                data["subtitle"], 14, WHITE)
    if data.get("badge"):
        badge = rect(slide, Inches(1.00), Inches(4.08), Inches(2.25), Inches(0.34), BLUE, WHITE)
        badge.fill.transparency = 20
        textbox(slide, Inches(1.12), Inches(4.17), Inches(2.0), Inches(0.12),
                data["badge"], 6.5, WHITE)
    rect(slide, Inches(8.30), Inches(1.20), Inches(3.10), Inches(0.012), WHITE)
    rect(slide, Inches(8.30), Inches(1.75), Inches(2.25), Inches(0.012), WHITE)
    rect(slide, Inches(8.30), Inches(2.30), Inches(3.45), Inches(0.012), WHITE)
    textbox(slide, Inches(11.35), Inches(6.82), Inches(1.05), Inches(0.12),
            str(page_no), 6, WHITE, align=PP_ALIGN.RIGHT)


def section_slide(prs: Presentation, data: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, prs.slide_width, prs.slide_height, BLUE)
    textbox(slide, Inches(2.2), Inches(3.05), Inches(4.55), Inches(0.62),
            data.get("title", "Section"), 31, WHITE, True, PP_ALIGN.RIGHT)
    rect(slide, Inches(7.0), Inches(3.42), Inches(2.55), Inches(0.012), WHITE)
    items = data.get("items") or []
    right_text = data.get("subtitle") or "\n".join(str(item) for item in items[:4])
    if right_text:
        textbox(slide, Inches(9.85), Inches(3.04), Inches(2.35), Inches(0.8),
                right_text, 12.5, WHITE, True)
    textbox(slide, Inches(11.35), Inches(6.82), Inches(1.05), Inches(0.12),
            str(page_no), 6, WHITE, align=PP_ALIGN.RIGHT)


def content_slide(prs: Presentation, data: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = data.get("title", "Untitled")
    body = data.get("body", "")
    add_heading(slide, title)
    textbox(slide, Inches(0.62), Inches(1.55), Inches(5.25), Inches(3.35), body, 13.5, TEXT)
    rect(slide, Inches(6.28), Inches(1.42), Inches(5.55), Inches(2.75), PALE_BLUE, CARD_BORDER, True)
    textbox(slide, Inches(6.58), Inches(1.78), Inches(1.0), Inches(0.25),
            f"{page_no:02d}", 14, BLUE, True)
    textbox(slide, Inches(7.08), Inches(1.76), Inches(4.25), Inches(0.32), title, 15, TEXT, True)
    rect(slide, Inches(6.58), Inches(2.34), Inches(4.9), Inches(0.01), CARD_BORDER)
    card_text = data.get("callout") or body
    textbox(slide, Inches(6.58), Inches(2.70), Inches(4.8), Inches(0.95), card_text, 12, TEXT)
    add_optional_progress(slide, page_no, data)


def message_slide(prs: Presentation, data: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    small_label(slide, Inches(0.65), Inches(0.60), data.get("eyebrow", "KEY MESSAGE"))
    textbox(slide, Inches(0.62), Inches(0.90), Inches(9.6), Inches(1.2),
            data.get("title", "Untitled"), 28, BLUE, True)
    if data.get("message"):
        textbox(slide, Inches(0.72), Inches(2.35), Inches(5.5), Inches(1.85),
                data["message"], 20, TEXT, True)
    if data.get("body"):
        textbox(slide, Inches(0.78), Inches(4.45), Inches(5.4), Inches(0.85),
                data["body"], 11.5, MID_GRAY)
    cards = data.get("cards") or []
    for idx, card in enumerate(cards[:3]):
        y = Inches(1.72 + idx * 1.30)
        fill = PALE_BLUE if data.get("active") == idx else WHITE
        rect(slide, Inches(7.05), y, Inches(4.65), Inches(0.95), fill, CARD_BORDER, True)
        textbox(slide, Inches(7.34), y + Inches(0.20), Inches(0.45), Inches(0.18),
                f"{idx + 1}", 13, BLUE, True)
        textbox(slide, Inches(7.88), y + Inches(0.18), Inches(3.45), Inches(0.2),
                card.get("title", ""), 12.5, TEXT, True)
        textbox(slide, Inches(7.88), y + Inches(0.50), Inches(3.55), Inches(0.18),
                card.get("body", ""), 8.5, MID_GRAY)
    add_optional_progress(slide, page_no, data)


def two_column_slide(prs: Presentation, data: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_heading(slide, data.get("title", "Untitled"))
    cols = data.get("columns") or []
    while len(cols) < 2:
        cols.append({})
    for idx, col in enumerate(cols[:2]):
        x = Inches(0.72 + idx * 5.85)
        fill = PALE_BLUE if data.get("active") == idx else WHITE
        rect(slide, x, Inches(1.65), Inches(5.05), Inches(3.5), fill, CARD_BORDER, True)
        textbox(slide, x + Inches(0.32), Inches(1.95), Inches(4.4), Inches(0.3),
                col.get("title", ""), 15, BLUE if data.get("active") == idx else TEXT, True)
        textbox(slide, x + Inches(0.32), Inches(2.45), Inches(4.35), Inches(2.05),
                col.get("body", ""), 12.5, TEXT)
    add_optional_progress(slide, page_no, data)


def compare_slide(prs: Presentation, data: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_heading(slide, data.get("title", "Untitled"))
    cols = data.get("columns") or []
    while len(cols) < 2:
        cols.append({})
    for idx, col in enumerate(cols[:2]):
        x = Inches(0.72 + idx * 5.88)
        fill = PALE_BLUE if data.get("active") == idx else WHITE
        rect(slide, x, Inches(1.58), Inches(5.1), Inches(3.95), fill, CARD_BORDER, True)
        rect(slide, x, Inches(1.58), Inches(5.1), Inches(0.09), BLUE_2)
        textbox(slide, x + Inches(0.35), Inches(1.92), Inches(4.35), Inches(0.30),
                col.get("title", ""), 17, BLUE if data.get("active") == idx else TEXT, True)
        textbox(slide, x + Inches(0.36), Inches(2.58), Inches(4.25), Inches(2.05),
                bullet_text(col.get("items") or col.get("body")), 12.0, TEXT)
    add_optional_progress(slide, page_no, data)


def process_slide(prs: Presentation, data: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_heading(slide, data.get("title", "Process"))
    steps = [str(step) for step in data.get("steps", ["Requirements", "Prototype", "Build"])]
    details = [str(detail) for detail in data.get("details", [])]
    active = data.get("active")
    active = int(active) if active is not None else None
    count = max(1, len(steps))
    width = Inches(11.3 / count)
    for idx, step in enumerate(steps):
        x = Inches(0.72) + idx * width
        fill = BLUE_2 if idx == active else PALE_BLUE
        text_color = WHITE if idx == active else TEXT
        rect(slide, x, Inches(2.40), width - Inches(0.08), Inches(1.35), fill, CARD_BORDER, True)
        textbox(slide, x + Inches(0.12), Inches(2.68), width - Inches(0.32), Inches(0.2),
                step, 12, text_color, True, PP_ALIGN.CENTER)
        if idx < len(details) and details[idx]:
            textbox(slide, x + Inches(0.18), Inches(3.15), width - Inches(0.42), Inches(0.35),
                    details[idx], 8.2, WHITE if idx == active else MID_GRAY, False, PP_ALIGN.CENTER)
    if data.get("body"):
        textbox(slide, Inches(0.74), Inches(4.25), Inches(10.8), Inches(0.7), data["body"], 12.5, TEXT)
    add_optional_progress(slide, page_no, data)


def roadmap_slide(prs: Presentation, data: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_heading(slide, data.get("title", "Roadmap"))
    phases = data.get("phases") or []
    if not phases:
        phases = [{"title": item, "body": ""} for item in data.get("steps", [])]
    active = data.get("active")
    active = int(active) if active is not None else None
    count = min(max(len(phases), 1), 4)
    rect(slide, Inches(0.95), Inches(3.08), Inches(10.9), Inches(0.035), LIGHT_GRAY)
    for idx, phase in enumerate(phases[:4]):
        x = Inches(0.85 + idx * (10.9 / count))
        color = BLUE_2 if idx == active else PALE_BLUE
        rect(slide, x, Inches(2.82), Inches(0.50), Inches(0.50), color, CARD_BORDER, True)
        textbox(slide, x + Inches(0.15), Inches(2.98), Inches(0.2), Inches(0.1),
                str(idx + 1), 9, WHITE if idx == active else BLUE, True,
                PP_ALIGN.CENTER)
        rect(slide, x - Inches(0.25), Inches(3.72), Inches(2.45), Inches(1.25),
             PALE_BLUE if idx == active else WHITE, CARD_BORDER, True)
        textbox(slide, x, Inches(3.98), Inches(1.85), Inches(0.25),
                phase.get("title", ""), 12.5, BLUE if idx == active else TEXT, True)
        textbox(slide, x, Inches(4.40), Inches(1.85), Inches(0.36),
                phase.get("body", ""), 8.5, MID_GRAY)
    if data.get("message"):
        textbox(slide, Inches(0.75), Inches(1.38), Inches(10.8), Inches(0.45),
                data["message"], 14.5, TEXT, True)
    add_optional_progress(slide, page_no, data)


def challenge_solution_slide(prs: Presentation, data: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_heading(slide, data.get("title", "Challenges and responses"))
    pairs = data.get("pairs") or []
    headers = [data.get("left_label", "課題"), data.get("right_label", "対策")]
    textbox(slide, Inches(0.95), Inches(1.47), Inches(4.9), Inches(0.24), headers[0], 12, BLUE, True)
    textbox(slide, Inches(6.65), Inches(1.47), Inches(4.9), Inches(0.24), headers[1], 12, BLUE, True)
    for idx, pair in enumerate(pairs[:5]):
        y = Inches(1.88 + idx * 0.78)
        rect(slide, Inches(0.82), y, Inches(5.10), Inches(0.58), WHITE, CARD_BORDER, True)
        rect(slide, Inches(6.45), y, Inches(5.10), Inches(0.58), WHITE, CARD_BORDER, True)
        textbox(slide, Inches(1.08), y + Inches(0.17), Inches(4.45), Inches(0.14),
                pair.get("challenge", ""), 9.2, TEXT, True)
        textbox(slide, Inches(6.72), y + Inches(0.15), Inches(4.38), Inches(0.18),
                pair.get("solution", ""), 9.2, TEXT)
        rect(slide, Inches(5.98), y + Inches(0.27), Inches(0.30), Inches(0.02), BLUE_2)
    add_optional_progress(slide, page_no, data)


def decision_slide(prs: Presentation, data: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_heading(slide, data.get("title", "Next decisions"))
    rect(slide, Inches(0.72), Inches(1.55), Inches(3.35), Inches(3.95), BLUE, None, True)
    textbox(slide, Inches(1.02), Inches(2.08), Inches(2.55), Inches(0.85),
            data.get("message", "Decide the first pilot."), 19, WHITE, True)
    textbox(slide, Inches(1.04), Inches(3.55), Inches(2.42), Inches(0.78),
            data.get("submessage", ""), 10.5, WHITE)
    items = data.get("items") or []
    for idx, item in enumerate(items[:5]):
        y = Inches(1.65 + idx * 0.72)
        rect(slide, Inches(4.70), y, Inches(6.65), Inches(0.52), WHITE, CARD_BORDER, True)
        textbox(slide, Inches(5.02), y + Inches(0.17), Inches(0.32), Inches(0.10),
                f"{idx + 1}", 8, BLUE, True)
        textbox(slide, Inches(5.45), y + Inches(0.14), Inches(5.35), Inches(0.14),
                item, 10.5, TEXT)
    add_optional_progress(slide, page_no, data)


def build_deck(spec: dict[str, Any], output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    clear_slides(prs)
    slides = spec.get("slides") or [
        {"type": "cover", "title": spec.get("title", "Untitled"), "subtitle": spec.get("subtitle", "")}
    ]
    builders = {
        "cover": cover_slide,
        "section": section_slide,
        "content": content_slide,
        "message": message_slide,
        "two-column": two_column_slide,
        "compare": compare_slide,
        "process": process_slide,
        "roadmap": roadmap_slide,
        "challenge-solution": challenge_solution_slide,
        "decision": decision_slide,
    }
    for page_no, slide_spec in enumerate(slides, 1):
        builder = builders.get(slide_spec.get("type", "content"), content_slide)
        builder(prs, slide_spec, page_no)
    prs.save(output)


def restyle(input_path: Path, output: Path) -> None:
    source = Presentation(input_path)
    extracted = [text_of_slide(slide) for slide in source.slides]
    slides: list[dict[str, Any]] = []
    for idx, texts in enumerate(extracted):
        title = texts[0] if texts else f"Slide {idx + 1}"
        body = "\n".join(texts[1:])
        if idx == 0:
            slides.append({"type": "cover", "title": title, "subtitle": body})
        elif idx == len(extracted) - 1 and len(extracted) > 2:
            slides.append({"type": "section", "title": title, "subtitle": body})
        else:
            slides.append({"type": "content", "title": title, "body": body})
    build_deck({"slides": slides}, output)


def sample(output: Path) -> None:
    build_deck(
        {
            "slides": [
                {
                    "type": "cover",
                    "title": "Robust Blue Presentation",
                    "subtitle": "Editable PowerPoint template",
                    "badge": "Draft",
                },
                {
                    "type": "section",
                    "title": "Requirements",
                    "items": ["Define purpose", "Understand constraints", "Prepare worksheet"],
                },
                {
                    "type": "message",
                    "title": "Define the purpose before designing the screen",
                    "message": "A good dashboard starts from the decision it enables.",
                    "body": "Clarify the user, decision, and action before choosing charts or layout.",
                    "cards": [
                        {"title": "User", "body": "Who acts on the information"},
                        {"title": "Decision", "body": "What needs to be decided"},
                        {"title": "Action", "body": "What changes after the decision"},
                    ],
                    "active": 0,
                },
                {
                    "type": "roadmap",
                    "title": "Creation flow",
                    "steps": ["Requirements", "Prototype", "Build"],
                    "message": "Move from purpose and constraints into layout exploration before implementation.",
                    "phases": [
                        {"title": "Requirements", "body": "Purpose and constraints"},
                        {"title": "Prototype", "body": "Structure and interaction"},
                        {"title": "Build", "body": "Implementation and checks"},
                    ],
                },
            ]
        },
        output,
    )


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    sub = parser.add_subparsers(dest="command", required=True)

    create = sub.add_parser("create")
    create.add_argument("--slides", required=True, type=Path, help="JSON slide specification.")
    create.add_argument("--output", required=True, type=Path)

    restyle_cmd = sub.add_parser("restyle")
    restyle_cmd.add_argument("--input", required=True, type=Path)
    restyle_cmd.add_argument("--output", required=True, type=Path)

    sample_cmd = sub.add_parser("sample")
    sample_cmd.add_argument("--output", required=True, type=Path)

    args = parser.parse_args()
    if args.command == "create":
        build_deck(json.loads(args.slides.read_text(encoding="utf-8")), args.output)
    elif args.command == "restyle":
        restyle(args.input, args.output)
    elif args.command == "sample":
        sample(args.output)


if __name__ == "__main__":
    main()
